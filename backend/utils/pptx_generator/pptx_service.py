"""
PPTX Generator Service
Core logic for generating PowerPoint presentations using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import os
import tempfile
import shutil
from typing import List, Dict, Any, Optional


class PPTXGeneratorService:
    """
    Service for generating PPTX files from templates
    """
    
    def __init__(self, templates_dir: str = None):
        """Initialize with templates directory"""
        if templates_dir is None:
            self.templates_dir = os.path.join(os.path.dirname(__file__), 'templates')
        else:
            self.templates_dir = templates_dir
            
        # Ensure templates directory exists
        os.makedirs(self.templates_dir, exist_ok=True)
    
    def get_available_templates(self) -> List[Dict[str, str]]:
        """List available template files"""
        templates = []
        
        if os.path.exists(self.templates_dir):
            for filename in os.listdir(self.templates_dir):
                if filename.endswith('.pptx'):
                    templates.append({
                        'name': filename.replace('.pptx', ''),
                        'path': os.path.join(self.templates_dir, filename),
                        'filename': filename
                    })
        
        return templates
    
    def generate_presentation(
        self,
        template_path: str,
        lesson_title: str,
        slides: List[Dict[str, Any]],
        title_bg_path: Optional[str] = None,
        content_bg_path: Optional[str] = None
    ) -> str:
        """
        Generate a PPTX presentation from template and content
        Layout matches pptx_creator.py exactly
        
        Args:
            template_path: Path to template PPTX file, or 'blank' for new presentation
            lesson_title: Title of the lesson/presentation
            slides: List of slide content dictionaries
            title_bg_path: Path to background image for title slide
            content_bg_path: Path to background image for content slides
            
        Returns:
            Path to generated PPTX file
        """
        # Log background paths for debugging
        print(f"[PPTX] title_bg_path: {title_bg_path}")
        print(f"[PPTX] content_bg_path: {content_bg_path}")
        
        # Create presentation (from template or blank)
        if template_path and os.path.exists(template_path) and template_path != 'blank':
            prs = Presentation(template_path)
        else:
            # Create blank presentation - MATCH pptx_creator.py exactly
            prs = Presentation()
            prs.slide_width = Inches(10)      # Match pptx_creator.py line 51
            prs.slide_height = Inches(5.625)  # Match pptx_creator.py line 52
        
        # Clear existing slides if using template (keep layouts)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # Add slides following pptx_creator.py structure:
        # - Slide 0: Title slide
        # - Slide 1: Agenda/Objectives (2-column)
        # - Slide 2+: Content slides
        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get('slideType', 'content')
            slide_index = slide_data.get('slideIndex', i)
            
            # Use layout 6 (blank) for ALL slides - match pptx_creator.py
            # Find true blank layout (no placeholders) to avoid "CLICK TO EDIT" text
            blank_layout = self._find_blank_layout(prs)
            slide = prs.slides.add_slide(blank_layout)
            
            # Clear any placeholder shapes that may exist from template layout
            self._clear_placeholders(slide)
            
            if slide_type == 'title' or slide_index == 0:
                self._add_title_slide(slide, slide_data, title_bg_path)
            elif slide_type in ['agenda', 'objectives'] or slide_index == 1:
                self._add_agenda_slide(slide, slide_data, content_bg_path)
            else:
                self._add_content_slide_v2(slide, slide_data, content_bg_path)
            
            # Add speaker notes
            speaker_note = slide_data.get('speakerNote', '')
            if speaker_note:
                self._add_speaker_notes(slide, speaker_note)
            
            # Add audio with auto-play
            audio_path = slide_data.get('audioPath')
            if audio_path and os.path.exists(audio_path):
                self._add_audio_with_autoplay(slide, audio_path)
        
        # Save to temp file
        output_path = tempfile.mktemp(suffix='.pptx')
        prs.save(output_path)
        
        return output_path
    
    def _find_blank_layout(self, prs: Presentation):
        """
        Find a truly blank layout (no placeholder shapes) from the presentation.
        Tries: layout 6 (standard blank), then searches for layout with no placeholders,
        finally falls back to last layout.
        """
        layouts = prs.slide_layouts
        
        # Try layout 6 first (standard blank in most templates)
        if len(layouts) > 6:
            layout_6 = layouts[6]
            # Check if it has placeholder shapes
            if len(list(layout_6.placeholders)) == 0:
                return layout_6
        
        # Search for a layout with no placeholders (truly blank)
        for layout in layouts:
            if len(list(layout.placeholders)) == 0:
                return layout
        
        # If no truly blank layout, find one with minimal placeholders
        min_placeholders = float('inf')
        best_layout = layouts[-1]
        for layout in layouts:
            num_ph = len(list(layout.placeholders))
            if num_ph < min_placeholders:
                min_placeholders = num_ph
                best_layout = layout
        
        return best_layout
    
    def _clear_placeholders(self, slide):
        """
        Remove all placeholder shapes from a slide.
        This prevents 'CLICK TO EDIT MASTER TITLE STYLE' text from appearing.
        Must iterate backward or collect shapes first because we're modifying the list.
        """
        # Collect all placeholder shapes first, then remove
        # Check for shapes with placeholder_format or specific text patterns
        shapes_to_remove = []
        
        for shape in slide.shapes:
            try:
                # Method 1: Check is_placeholder attribute
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    shapes_to_remove.append(shape)
                    continue
                
                # Method 2: Check if shape has placeholder format
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                    shapes_to_remove.append(shape)
                    continue
                    
                # Method 3: Check for placeholder-like text (CLICK TO EDIT, Click to add)
                if hasattr(shape, 'text_frame'):
                    text = shape.text_frame.text.lower() if shape.text_frame.text else ''
                    if 'click to' in text or 'master' in text:
                        shapes_to_remove.append(shape)
                        continue
            except Exception:
                # Skip shapes that cause errors when checking
                pass
        
        # Remove all identified placeholder shapes
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                if sp.getparent() is not None:
                    sp.getparent().remove(sp)
                    print(f"[PPTX] Removed placeholder shape")
            except Exception as e:
                print(f"[PPTX] Could not remove placeholder: {e}")
    
    def _add_title_slide(self, slide, slide_data: Dict[str, Any], bg_path: Optional[str] = None):
        """
        Add Title Slide - matches pptx_creator.py lines 54-80
        Custom textbox for title/subtitle, white text on dark background
        """
        # Add background image if provided (MUST be added FIRST so it's behind text)
        if bg_path and os.path.exists(bg_path):
            try:
                # Fill entire slide with background image
                bg_pic = slide.shapes.add_picture(
                    bg_path,
                    Inches(0), Inches(0),
                    width=Inches(10), height=Inches(5.625)
                )
                # Send to back
                spTree = slide.shapes._spTree
                spTree.insert(2, bg_pic._element)  # Insert after shape tree start
                print(f"[PPTX] Title slide background added: {bg_path}")
            except Exception as e:
                print(f"[PPTX] Could not add title background: {e}")
        
        title = slide_data.get('title', '')
        content = slide_data.get('content', [])
        subtitle = content[0] if content else ''
        
        # Title textbox - centered at position (1, 2)
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = title_shape.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle textbox
        if subtitle:
            subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            tf_sub = subtitle_shape.text_frame
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.name = "Arial"
            p_sub.font.size = Pt(22)
            p_sub.font.color.rgb = RGBColor(255, 255, 255)  # White text
            p_sub.alignment = PP_ALIGN.CENTER
    
    def _add_slide(
        self,
        prs: Presentation,
        layouts,
        slide_data: Dict[str, Any]
    ):
        """Add a single slide to the presentation"""
        
        slide_type = slide_data.get('slideType', 'content')
        title = slide_data.get('title', '')
        content = slide_data.get('content', [])
        bullets = slide_data.get('bullets')  # Structured bullets from AI
        image_path = slide_data.get('imagePath')
        audio_path = slide_data.get('audioPath')
        speaker_note = slide_data.get('speakerNote', '')
        slide_index = slide_data.get('slideIndex', 0)
        
        # Select appropriate layout
        layout_idx = self._get_layout_index(layouts, slide_type)
        slide_layout = layouts[layout_idx]
        
        # Add slide
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Add content based on slide type
        if slide_type == 'title':
            self._add_title_slide_content(slide, title, content)
        elif slide_type in ['agenda', 'objectives'] or slide_index == 1:
            # Agenda/objectives slide with 2-column layout (matching pptx_creator.py)
            self._add_agenda_slide(slide, title, content, bullets)
        else:
            self._add_content_slide(slide, title, content, image_path, bullets)
        
        # Add speaker notes
        if speaker_note:
            self._add_speaker_notes(slide, speaker_note)
        
        # Add audio with auto-play
        if audio_path and os.path.exists(audio_path):
            self._add_audio_with_autoplay(slide, audio_path)
    
    def _add_agenda_slide(self, slide, slide_data: Dict[str, Any], bg_path: Optional[str] = None):
        """
        Add Agenda/Objectives slide with 2-column layout
        Matches pptx_creator.py lines 89-130
        """
        # Add background image if provided (MUST be added FIRST so it's behind text)
        if bg_path and os.path.exists(bg_path):
            try:
                bg_pic = slide.shapes.add_picture(
                    bg_path,
                    Inches(0), Inches(0),
                    width=Inches(10), height=Inches(5.625)
                )
                # Send to back
                spTree = slide.shapes._spTree
                spTree.insert(2, bg_pic._element)
                print(f"[PPTX] Agenda slide background added: {bg_path}")
            except Exception as e:
                print(f"[PPTX] Could not add agenda background: {e}")
        
        title = slide_data.get('title', 'Nội dung bài học')
        bullets = slide_data.get('bullets', [])
        content = slide_data.get('content', [])
        
        # Title textbox - white text at top
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        p_title = title_shape.text_frame.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p_title.alignment = PP_ALIGN.CENTER
        
        # Use bullets if available, else fallback to content
        items = []
        if bullets:
            items = [f"{b.get('emoji', '')} {b.get('point', '') or b.get('description', '')}" for b in bullets]
        elif content:
            items = content if isinstance(content, list) else [content]
        
        if not items:
            return
        
        mid_point = (len(items) + 1) // 2
        
        # Left Column
        left_col = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4.0))
        tf_left = left_col.text_frame
        tf_left.clear()
        tf_left.word_wrap = True
        for item in items[:mid_point]:
            p = tf_left.add_paragraph()
            p.text = item
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(58, 102, 77)
            p.space_after = Pt(16)
        
        # Right Column
        right_col = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.5), Inches(4.0))
        tf_right = right_col.text_frame
        tf_right.clear()
        tf_right.word_wrap = True
        for item in items[mid_point:]:
            p = tf_right.add_paragraph()
            p.text = item
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(58, 102, 77)
            p.space_after = Pt(16)
    
    def _add_content_slide_v2(self, slide, slide_data: Dict[str, Any], bg_path: Optional[str] = None):
        """
        Add content slide with exact pptx_creator.py layout (lines 132-177)
        - Title: white text at (0.5, 0.2)
        - If image: content at left (4.5"), image at right (5.5, 1.5)
        - If no image: content full width (9")
        - Bullets: emoji+point bold 22pt, description 18pt
        """
        # Add background image if provided (MUST be added FIRST so it's behind text)
        if bg_path and os.path.exists(bg_path):
            try:
                bg_pic = slide.shapes.add_picture(
                    bg_path,
                    Inches(0), Inches(0),
                    width=Inches(10), height=Inches(5.625)
                )
                # Send to back
                spTree = slide.shapes._spTree
                spTree.insert(2, bg_pic._element)
                print(f"[PPTX] Content slide background added: {bg_path}")
            except Exception as e:
                print(f"[PPTX] Could not add content background: {e}")
        
        title = slide_data.get('title', '')
        bullets = slide_data.get('bullets', [])
        content = slide_data.get('content', [])
        image_path = slide_data.get('imagePath')
        
        # Check if image exists
        has_image = image_path and os.path.exists(image_path)
        if image_path and not has_image:
            print(f"[PPTX] Image NOT found: {image_path}")
        elif has_image:
            print(f"[PPTX] Image found: {image_path}")
        
        # Title textbox - white text at top (match pptx_creator.py lines 138-144)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        p_title = title_shape.text_frame.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.color.rgb = RGBColor(255, 255, 255)  # White text
        p_title.alignment = PP_ALIGN.CENTER
        
        # Content textbox - width depends on image (match pptx_creator.py lines 146-153)
        if has_image:
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4.0))
        else:
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.0))
        
        content_frame = content_shape.text_frame
        content_frame.clear()
        content_frame.word_wrap = True
        
        # Add bullets (match pptx_creator.py lines 159-177)
        if bullets:
            for bullet in bullets:
                emoji = bullet.get('emoji', '')
                point = bullet.get('point', '')
                description = bullet.get('description', '')
                
                if point:
                    # Bold emoji+point (22pt)
                    p_point = content_frame.add_paragraph()
                    p_point.text = f'{emoji} {point}' if emoji else point
                    p_point.font.name = "Arial"
                    p_point.font.bold = True
                    p_point.font.size = Pt(22)
                    p_point.font.color.rgb = RGBColor(58, 102, 77)
                    p_point.space_after = Pt(2)
                    
                    # Description (18pt, indented)
                    if description:
                        p_desc = content_frame.add_paragraph()
                        p_desc.text = description
                        p_desc.font.name = "Arial"
                        p_desc.font.size = Pt(18)
                        p_desc.font.color.rgb = RGBColor(58, 102, 77)
                        p_desc.space_before = Pt(0)
                        p_desc.space_after = Pt(8)
                        p_desc.level = 1
                else:
                    # Golden Rule/Definition: Just description (20pt)
                    p_point = content_frame.add_paragraph()
                    p_point.text = description
                    p_point.font.name = "Arial"
                    p_point.font.size = Pt(20)
                    p_point.font.color.rgb = RGBColor(58, 102, 77)
                    p_point.space_after = Pt(8)
        
        elif content:
            # Fallback to flat content array
            for item in content:
                p = content_frame.add_paragraph()
                p.text = f"• {item}" if not item.startswith('•') else item
                p.font.name = "Arial"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(58, 102, 77)
                p.space_after = Pt(6)
        
        # Add image if exists (match pptx_creator.py line 149)
        if has_image:
            try:
                slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), height=Inches(3))
            except Exception as e:
                print(f"[PPTX] Could not add picture {image_path}: {e}")
    
    def _get_layout_index(self, layouts, slide_type: str) -> int:
        """Get appropriate layout index for slide type"""
        layout_map = {
            'title': 0,        # Title Slide
            'objectives': 1,   # Title and Content
            'content': 1,      # Title and Content
            'summary': 1,      # Title and Content
            'blank': 6,        # Blank
        }
        
        idx = layout_map.get(slide_type, 1)
        
        # Ensure index is valid
        if idx >= len(layouts):
            idx = min(1, len(layouts) - 1)
        
        return idx
    
    def _add_title_slide_content(
        self,
        slide,
        title: str,
        content: List[str]
    ):
        """Add content to title slide"""
        # Title is already set, add subtitle if content exists
        if content and len(content) > 0:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape != slide.shapes.title:
                    if hasattr(shape, 'text_frame'):
                        shape.text = content[0]
                        break
    
    def _add_content_slide(
        self,
        slide,
        title: str,
        content: List[str],
        image_path: Optional[str],
        bullets: Optional[List[Dict[str, str]]] = None
    ):
        """
        Add content to regular content slide with proper layout.
        
        Format matches pptx_creator.py:
        - If bullet has point: Bold emoji+point (22pt), then description (18pt) on next line
        - If bullet has no point (definition/concept): Just description (20pt)
        """
        
        # Check image path and log for debugging
        has_image = False
        if image_path:
            if os.path.exists(image_path):
                has_image = True
                print(f"[PPTX] Image found: {image_path}")
            else:
                print(f"[PPTX] Image NOT found: {image_path}")
        
        # Find content placeholder or create text box
        content_placeholder = None
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                content_placeholder = shape
                break
        
        # Create text frame for content (similar to pptx_creator.py)
        if content_placeholder:
            tf = content_placeholder.text_frame
            tf.clear()
            tf.word_wrap = True
        else:
            # Create text box if no placeholder
            left = Inches(0.5)
            top = Inches(1.2)  # Match pptx_creator.py
            
            if has_image:
                width = Inches(4.5)  # Make room for image on right
            else:
                width = Inches(9)  # Full width match pptx_creator.py
            
            height = Inches(4.0)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
        
        # Use structured bullets if available, otherwise fallback to flat content
        if bullets:
            # Limit bullets to prevent overflow (max 5 for structured format)
            MAX_BULLETS = 5
            if len(bullets) > MAX_BULLETS:
                bullets = bullets[:MAX_BULLETS]
                print(f"Warning: Bullets truncated to {MAX_BULLETS}")
            
            first_para = True
            for bullet in bullets:
                emoji = bullet.get('emoji', '')
                point = bullet.get('point', '')
                description = bullet.get('description', '')
                
                if point:
                    # Format: Bold emoji+point, then description on next line
                    # (Matching pptx_creator.py lines 160-173)
                    p_point = tf.add_paragraph() if not first_para else tf.paragraphs[0]
                    p_point.text = f'{emoji} {point}' if emoji else point
                    p_point.font.name = "Arial"
                    p_point.font.bold = True
                    p_point.font.size = Pt(22)
                    p_point.font.color.rgb = RGBColor(58, 102, 77)  # Green color
                    p_point.space_after = Pt(2)
                    
                    if description:
                        p_desc = tf.add_paragraph()
                        p_desc.text = description
                        p_desc.font.name = "Arial"
                        p_desc.font.size = Pt(18)
                        p_desc.font.color.rgb = RGBColor(58, 102, 77)
                        p_desc.space_before = Pt(0)
                        p_desc.space_after = Pt(8)
                        p_desc.level = 1
                else:
                    # Golden Rule: Definition/concept - just show full description
                    # (Matching pptx_creator.py lines 174-177)
                    p_point = tf.add_paragraph() if not first_para else tf.paragraphs[0]
                    p_point.text = description
                    p_point.font.name = "Arial"
                    p_point.font.size = Pt(20)
                    p_point.font.color.rgb = RGBColor(58, 102, 77)
                    p_point.space_after = Pt(8)
                
                first_para = False
        
        elif content:
            # Fallback: flat content array
            MAX_BULLETS = 6
            if len(content) > MAX_BULLETS:
                content = content[:MAX_BULLETS]
            
            # Calculate font size based on content length
            if len(content) <= 3:
                font_size = Pt(16)
            elif len(content) <= 5:
                font_size = Pt(14)
            else:
                font_size = Pt(12)
            
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {item}"
                p.font.name = "Arial"
                p.font.size = font_size
                p.font.color.rgb = RGBColor(58, 102, 77)
                p.space_after = Pt(6)
        
        # Add image if exists
        if has_image:
            self._add_image(slide, image_path, has_content=bool(content or bullets))
    
    def _add_image(self, slide, image_path: str, has_content: bool = False):
        """Add image to slide"""
        try:
            if has_content:
                # Image on right side
                left = Inches(7)
                top = Inches(1.5)
                width = Inches(5.5)
            else:
                # Image centered
                left = Inches(1.5)
                top = Inches(1.5)
                width = Inches(10)
            
            slide.shapes.add_picture(image_path, left, top, width=width)
        except Exception as e:
            print(f"Warning: Could not add image: {e}")
    
    def _add_speaker_notes(self, slide, notes: str):
        """Add speaker notes to slide"""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    def _add_audio_with_autoplay(self, slide, audio_path: str):
        """
        Add audio to slide with auto-play on slide transition.
        Audio icon is placed OUTSIDE the visible slide area (top-left corner)
        so students don't see it during presentation.
        """
        try:
            # Get file extension
            ext = os.path.splitext(audio_path)[1].lower()
            
            # Supported audio formats
            if ext not in ['.mp3', '.wav', '.m4a', '.wma']:
                print(f"Warning: Unsupported audio format: {ext}")
                return
            
            # Position audio icon OUTSIDE visible slide area
            # Negative position places it off-screen to the top-left
            left = Inches(-0.6)   # Off-screen to the left
            top = Inches(-0.6)    # Off-screen to the top
            width = Inches(0.4)
            height = Inches(0.4)
            
            # Add media shape
            movie = slide.shapes.add_movie(
                audio_path,
                left, top, width, height,
                mime_type=self._get_audio_mime_type(ext)
            )
            
            print(f"[PPTX] Audio added (hidden): {audio_path}")
            
            # Set auto-play using XML manipulation
            self._set_audio_autoplay(movie, slide)
            
        except Exception as e:
            print(f"Warning: Could not add audio: {e}")
    
    def _get_audio_mime_type(self, ext: str) -> str:
        """Get MIME type for audio format"""
        mime_map = {
            '.mp3': 'audio/mpeg',
            '.wav': 'audio/wav',
            '.m4a': 'audio/mp4',
            '.wma': 'audio/x-ms-wma',
        }
        return mime_map.get(ext, 'audio/mpeg')
    
    def _set_audio_autoplay(self, movie_shape, slide):
        """
        Set audio to auto-play when slide appears
        This requires XML manipulation of the slide timing
        """
        try:
            # Get the shape id
            shape_id = movie_shape.shape_id
            
            # Create timing XML for auto-play
            # This sets the audio to play automatically when the slide loads
            timing_xml = f'''
            <p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
                <p:tnLst>
                    <p:par>
                        <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
                            <p:childTnLst>
                                <p:seq concurrent="1" nextAc="seek">
                                    <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                                        <p:childTnLst>
                                            <p:par>
                                                <p:cTn id="3" fill="hold">
                                                    <p:stCondLst>
                                                        <p:cond delay="0"/>
                                                    </p:stCondLst>
                                                    <p:childTnLst>
                                                        <p:par>
                                                            <p:cTn id="4" fill="hold">
                                                                <p:stCondLst>
                                                                    <p:cond delay="0"/>
                                                                </p:stCondLst>
                                                                <p:childTnLst>
                                                                    <p:par>
                                                                        <p:cTn id="5" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" nodeType="afterEffect">
                                                                            <p:stCondLst>
                                                                                <p:cond delay="0"/>
                                                                            </p:stCondLst>
                                                                            <p:childTnLst>
                                                                                <p:cmd type="call" cmd="playFrom(0.0)">
                                                                                    <p:cBhvr>
                                                                                        <p:cTn id="6" dur="1" fill="hold"/>
                                                                                        <p:tgtEl>
                                                                                            <p:spTgt spid="{shape_id}"/>
                                                                                        </p:tgtEl>
                                                                                    </p:cBhvr>
                                                                                </p:cmd>
                                                                            </p:childTnLst>
                                                                        </p:cTn>
                                                                    </p:par>
                                                                </p:childTnLst>
                                                            </p:cTn>
                                                        </p:par>
                                                    </p:childTnLst>
                                                </p:cTn>
                                            </p:par>
                                        </p:childTnLst>
                                    </p:cTn>
                                </p:seq>
                            </p:childTnLst>
                        </p:cTn>
                    </p:par>
                </p:tnLst>
            </p:timing>
            '''
            
            # Parse and add timing element
            timing_elem = parse_xml(timing_xml)
            
            # Find existing timing or add new
            slide_elem = slide._element
            existing_timing = slide_elem.find(qn('p:timing'))
            
            if existing_timing is not None:
                slide_elem.remove(existing_timing)
            
            slide_elem.append(timing_elem)
            
        except Exception as e:
            print(f"Warning: Could not set auto-play for audio: {e}")


# Standalone test
if __name__ == "__main__":
    service = PPTXGeneratorService()
    
    # Test data
    test_slides = [
        {
            "slideIndex": 0,
            "slideType": "title",
            "title": "Bài 01: Giới thiệu Python",
            "content": ["Khóa học lập trình cơ bản"],
            "speakerNote": "Chào mừng các bạn đến với bài học đầu tiên."
        },
        {
            "slideIndex": 1,
            "slideType": "objectives",
            "title": "Mục tiêu bài học",
            "content": [
                "Hiểu Python là gì",
                "Cài đặt Python và IDE",
                "Viết chương trình Hello World"
            ],
            "speakerNote": "Sau bài học này, các bạn sẽ biết cách cài đặt và chạy Python."
        },
        {
            "slideIndex": 2,
            "slideType": "content",
            "title": "Python là gì?",
            "content": [
                "Ngôn ngữ lập trình bậc cao",
                "Dễ đọc, dễ học",
                "Phổ biến trong AI, Data Science, Web"
            ],
            "speakerNote": "Python được tạo ra bởi Guido van Rossum năm 1991."
        }
    ]
    
    output = service.generate_presentation(
        template_path='blank',
        lesson_title='Test Presentation',
        slides=test_slides
    )
    
    print(f"Generated: {output}")
