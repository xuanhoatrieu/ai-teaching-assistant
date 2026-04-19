"""
PPTX Generator Service
FastAPI server for generating PowerPoint presentations from templates
"""

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional, Dict
import tempfile
import os
import logging

logger = logging.getLogger("pptx-generator")

from pptx_service import PPTXGeneratorService

# VietNormalizer for TTS text normalization (local clone)
import sys
from pathlib import Path
_VIETNORMALIZER_PATH = str(Path(__file__).parent / "vietnormalizer")
if _VIETNORMALIZER_PATH not in sys.path:
    sys.path.insert(0, _VIETNORMALIZER_PATH)

try:
    from vietnormalizer import VietnameseNormalizer
    _default_normalizer = VietnameseNormalizer()
    logger.info(f"VietNormalizer loaded from local clone: {_VIETNORMALIZER_PATH}")
except ImportError as e:
    _default_normalizer = None
    logger.warning(f"vietnormalizer failed to load: {e}")

app = FastAPI(
    title="PPTX Generator Service",
    description="Generate PowerPoint presentations from templates with content, images, and audio",
    version="1.0.0"
)

# CORS for NestJS backend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3001", "http://localhost:5174"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize service
pptx_service = PPTXGeneratorService()


# ========================
# Text Normalization Models
# ========================

class DictEntry(BaseModel):
    original: str
    replacement: str

class NormalizeRequest(BaseModel):
    text: str
    enable_transliteration: bool = True
    custom_acronyms: Optional[List[DictEntry]] = None
    custom_words: Optional[List[DictEntry]] = None

class NormalizeResponse(BaseModel):
    normalized_text: str
    original_text: str
    changes_made: bool


class SlideContent(BaseModel):
    slideIndex: int
    title: str
    content: List[str] = []
    bullets: Optional[List[dict]] = None  # Structured bullets from AI
    imagePath: Optional[str] = None
    audioPath: Optional[str] = None
    speakerNote: Optional[str] = None
    slideType: Optional[str] = "content"


class GeneratePPTXRequest(BaseModel):
    templatePath: str
    lessonTitle: str
    slides: List[SlideContent]
    titleBgPath: Optional[str] = None     # Background image for title slide
    contentBgPath: Optional[str] = None   # Background image for content slides


@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "ok",
        "service": "pptx-generator",
        "normalizer_available": _default_normalizer is not None,
    }


# In-memory dictionary cache (loaded from DB via /reload-dictionaries)
_cached_acronyms: dict = {}
_cached_words: dict = {}


class ReloadDictionariesRequest(BaseModel):
    acronyms: List[DictEntry] = []
    words: List[DictEntry] = []


@app.post("/reload-dictionaries")
async def reload_dictionaries(request: ReloadDictionariesRequest):
    """
    Reload dictionaries from DB (called by NestJS on startup or when admin changes).
    Replaces all cached dictionaries in memory.
    """
    global _cached_acronyms, _cached_words

    _cached_acronyms = {e.original.lower(): e.replacement for e in request.acronyms}
    _cached_words = {e.original.lower(): e.replacement for e in request.words}

    logger.info(f"Dictionaries reloaded: {len(_cached_acronyms)} acronyms, {len(_cached_words)} words")
    return {
        "status": "ok",
        "acronyms_count": len(_cached_acronyms),
        "words_count": len(_cached_words),
    }


@app.post("/normalize", response_model=NormalizeResponse)
async def normalize_text(request: NormalizeRequest):
    """
    Normalize Vietnamese text for TTS.
    Uses cached DB dictionaries + optional per-request user overrides.
    """
    if _default_normalizer is None:
        raise HTTPException(
            status_code=503,
            detail="vietnormalizer is not installed"
        )

    try:
        normalizer = _default_normalizer

        # Build merged dicts: cached system + per-request user overrides
        merged_acronyms = dict(_cached_acronyms)
        merged_words = dict(_cached_words)

        if request.custom_acronyms:
            for entry in request.custom_acronyms:
                merged_acronyms[entry.original.lower()] = entry.replacement

        if request.custom_words:
            for entry in request.custom_words:
                merged_words[entry.original.lower()] = entry.replacement

        # Temporarily swap dicts on normalizer
        original_acronym_map = normalizer.acronym_map
        original_non_viet_map = normalizer.non_vietnamese_map
        original_replacements = normalizer.replacements

        try:
            normalizer.acronym_map = merged_acronyms
            normalizer.non_vietnamese_map = merged_words
            normalizer.replacements = {k: v for k, v in merged_words.items()}

            normalized = normalizer.normalize(
                request.text,
                enable_transliteration=request.enable_transliteration,
            )
        finally:
            normalizer.acronym_map = original_acronym_map
            normalizer.non_vietnamese_map = original_non_viet_map
            normalizer.replacements = original_replacements

        return NormalizeResponse(
            normalized_text=normalized,
            original_text=request.text,
            changes_made=normalized != request.text,
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Normalization failed: {str(e)}")


@app.get("/templates")
async def list_templates():
    """List available system templates"""
    templates = pptx_service.get_available_templates()
    return {"templates": templates}


@app.post("/generate")
async def generate_pptx(request: GeneratePPTXRequest):
    """
    Generate a PPTX file from template and content
    Returns the generated PPTX file
    """
    try:
        # Generate PPTX
        output_path = pptx_service.generate_presentation(
            template_path=request.templatePath,
            lesson_title=request.lessonTitle,
            slides=[s.model_dump() for s in request.slides],
            title_bg_path=request.titleBgPath,
            content_bg_path=request.contentBgPath
        )
        
        # Return the file
        return FileResponse(
            path=output_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{request.lessonTitle}.pptx"
        )
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX: {str(e)}")


@app.post("/generate-buffer")
async def generate_pptx_buffer(request: GeneratePPTXRequest):
    """
    Generate a PPTX file and return as base64 buffer
    Alternative to /generate for easier integration
    """
    import base64
    
    try:
        output_path = pptx_service.generate_presentation(
            template_path=request.templatePath,
            lesson_title=request.lessonTitle,
            slides=[s.model_dump() for s in request.slides],
            title_bg_path=request.titleBgPath,
            content_bg_path=request.contentBgPath
        )
        
        with open(output_path, 'rb') as f:
            buffer = base64.b64encode(f.read()).decode('utf-8')
        
        # Cleanup temp file
        os.unlink(output_path)
        
        return {
            "buffer": buffer,
            "filename": f"{request.lessonTitle}.pptx"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate PPTX: {str(e)}")


# ========================
# PPTX Audio Tool Endpoints
# ========================

import re
import base64
from pptx import Presentation as PptxPresentation


# Vietnamese character detection for bilingual note splitting
VN_CHARS = set("àáạảãâầấậẩẫăằắặẳẵèéẹẻẽêềếệểễìíịỉĩòóọỏõôồốộổỗơờớợởỡùúụủũưừứựửữỳýỵỷỹđ"
               "ÀÁẠẢÃÂẦẤẬẨẪĂẰẮẶẲẴÈÉẸẺẼÊỀẾỆỂỄÌÍỊỈĨÒÓỌỎÕÔỒỐỘỔỖƠỜỚỢỞỠÙÚỤỦŨƯỪỨỰỬỮỲÝỴỶỸĐ")


def _detect_vietnamese(text: str) -> bool:
    """Check if text contains Vietnamese-specific characters."""
    return any(c in VN_CHARS for c in text)


def _split_bilingual_notes(note_text: str) -> dict:
    """
    Split bilingual speaker notes.
    Separator: blank line (\\n\\n)
    Detection: Vietnamese chars → VN part, ASCII-only → EN part
    """
    if not note_text or not note_text.strip():
        return {"full": "", "en": "", "vi": "", "has_dual": False}

    # Split by double newline
    parts = [p.strip() for p in note_text.strip().split('\n\n') if p.strip()]

    if len(parts) < 2:
        # Single block — assign to both
        return {
            "full": note_text.strip(),
            "en": note_text.strip(),
            "vi": note_text.strip(),
            "has_dual": False,
        }

    # Try to detect which is Vietnamese and which is English
    part1_is_vn = _detect_vietnamese(parts[0])
    part2_is_vn = _detect_vietnamese(parts[1])

    if part1_is_vn and not part2_is_vn:
        return {
            "full": note_text.strip(),
            "vi": parts[0],
            "en": parts[1],
            "has_dual": True,
        }
    elif part2_is_vn and not part1_is_vn:
        return {
            "full": note_text.strip(),
            "en": parts[0],
            "vi": parts[1],
            "has_dual": True,
        }
    else:
        # Both same language or can't determine
        return {
            "full": note_text.strip(),
            "en": note_text.strip(),
            "vi": note_text.strip(),
            "has_dual": False,
        }


def _extract_slides_from_pptx(prs) -> list:
    slides_data = []

    for i, slide in enumerate(prs.slides):
        # Extract title (first text shape or placeholder title)
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()

        if not title:
            # Fallback: find first text shape
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip()
                    break

        # Extract content (all non-title text shapes)
        slide_content = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Skip if this is the title shape
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text and text != title:
                    slide_content.append(text)

        # Extract speaker notes
        note_full = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                note_full = notes_frame.text.strip()

        # Split bilingual notes
        note_data = _split_bilingual_notes(note_full)

        slides_data.append({
            "index": i,
            "title": title or f"Slide {i + 1}",
            "content": slide_content,
            "noteFull": note_data["full"],
            "noteEN": note_data["en"],
            "noteVN": note_data["vi"],
            "hasDual": note_data["has_dual"],
        })
    return slides_data


@app.post("/parse-pptx")
async def parse_pptx(file: UploadFile = File(...)):
    """
    Parse uploaded PPTX → extract slides with speaker notes + content.
    Supports bilingual notes (EN/VN separated by blank line).
    """
    try:
        # Save uploaded file to temp
        tmp_path = tempfile.mktemp(suffix='.pptx')
        content = await file.read()
        with open(tmp_path, 'wb') as f:
            f.write(content)

        prs = PptxPresentation(tmp_path)
        slides_data = _extract_slides_from_pptx(prs)

        # Cleanup temp file
        os.unlink(tmp_path)

        logger.info(f"Parsed PPTX: {len(slides_data)} slides from {file.filename}")
        return {"slides": slides_data, "totalSlides": len(slides_data)}

    except Exception as e:
        logger.error(f"Failed to parse PPTX: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse PPTX: {str(e)}")


class ParseLocalRequest(BaseModel):
    file_path: str


@app.post("/parse-pptx-local")
async def parse_pptx_local(request: ParseLocalRequest):
    """
    Parse PPTX directly from a local shared volume path.
    Saves massive memory and network overhead on VPS.
    """
    try:
        if not os.path.exists(request.file_path):
            raise HTTPException(status_code=404, detail=f"File not found: {request.file_path}")

        prs = PptxPresentation(request.file_path)
        slides_data = _extract_slides_from_pptx(prs)

        logger.info(f"Parsed Local PPTX: {len(slides_data)} slides from {request.file_path}")
        return {"slides": slides_data, "totalSlides": len(slides_data)}

    except Exception as e:
        logger.error(f"Failed to parse local PPTX: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse local PPTX: {str(e)}")


class InjectAudioRequest(BaseModel):
    pptxPath: str
    audioFiles: List[Dict]  # [{slideIndex: 0, audioPath: "/path/to/audio.wav"}]


@app.post("/inject-audio")
async def inject_audio(request: InjectAudioRequest):
    """
    Open original PPTX, inject audio into each slide with auto-play.
    Uses same _add_audio_with_autoplay() method from pptx_service.
    """
    try:
        if not os.path.exists(request.pptxPath):
            raise HTTPException(status_code=404, detail=f"PPTX file not found: {request.pptxPath}")

        prs = PptxPresentation(request.pptxPath)

        injected_count = 0
        for audio_info in request.audioFiles:
            slide_idx = audio_info.get("slideIndex", -1)
            audio_path = audio_info.get("audioPath", "")

            if slide_idx < 0 or slide_idx >= len(prs.slides):
                logger.warning(f"Invalid slide index: {slide_idx}")
                continue

            if not audio_path or not os.path.exists(audio_path):
                logger.warning(f"Audio not found: {audio_path}")
                continue

            slide = prs.slides[slide_idx]
            pptx_service._add_audio_with_autoplay(slide, audio_path)
            injected_count += 1
            logger.info(f"Injected audio into slide {slide_idx}: {audio_path}")

        # Save to temp file
        output_path = tempfile.mktemp(suffix='.pptx')
        prs.save(output_path)

        # Read and encode as base64
        with open(output_path, 'rb') as f:
            buffer = base64.b64encode(f.read()).decode('utf-8')

        # Cleanup
        os.unlink(output_path)

        logger.info(f"Injected audio into {injected_count} slides")
        return {"buffer": buffer, "injectedCount": injected_count}

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to inject audio: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to inject audio: {str(e)}")


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=3002)
